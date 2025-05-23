import os
import time
import logging
from typing import List, Dict, Any
import pypdf
import docx2txt
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
import shutil
import torch

# Import GPU utilities
from app.core.gpu_utils import device_config

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Directories
DOCUMENTS_DIR = "data/documents"
EMBEDDINGS_DIR = "data/embeddings"

# Ensure directories exist
os.makedirs(DOCUMENTS_DIR, exist_ok=True)
os.makedirs(EMBEDDINGS_DIR, exist_ok=True)

# Initialize embeddings model with GPU support
# Force CUDA for embeddings if available
embeddings_device = "cuda" if torch.cuda.is_available() else "cpu"
logger.info(f"Initializing embeddings model on {embeddings_device} device")

# Use larger batch size for GPU to improve performance
batch_size = 128 if embeddings_device == "cuda" else 32  # Increased from 64 to 128 for faster processing

embeddings = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2",
    model_kwargs={"device": embeddings_device},
    encode_kwargs={
        "normalize_embeddings": True, 
        "device": embeddings_device,
        "batch_size": batch_size  # Larger batch size for GPU processing
    },
    show_progress=False  # Disable progress bar for faster processing
)

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    text = ""
    try:
        with open(file_path, "rb") as file:
            pdf = pypdf.PdfReader(file)
            for page in pdf.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path}: {e}")
    return text

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file"""
    try:
        text = docx2txt.process(file_path)
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {e}")
    return text

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from TXT file"""
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()
    except UnicodeDecodeError:
        # Try with a different encoding if UTF-8 fails
        try:
            with open(file_path, "r", encoding="latin-1") as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error reading TXT file {file_path}: {e}")
            return ""
    except Exception as e:
        logger.error(f"Error reading TXT file {file_path}: {e}")
        return ""

def extract_text(file_path: str) -> str:
    """Extract text from a document based on its extension"""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    else:
        logger.warning(f"Unsupported file format: {ext}")
        return ""

def create_document_chunks(text: str, filename: str) -> List[Dict[str, Any]]:
    """Split document text into chunks using a dynamic, content-aware strategy"""
    # Get file extension and size
    _, ext = os.path.splitext(filename)
    ext = ext.lower()
    text_size = len(text)
    
    # Detect document type and content characteristics
    is_markdown = ext == '.md' or (ext == '.txt' and ('#' in text or '##' in text))
    is_code = ext in ('.py', '.js', '.java', '.cpp', '.c', '.h', '.cs', '.php', '.ts', '.go', '.rb')
    has_long_paragraphs = max([len(p) for p in text.split('\n\n') if p.strip()], default=0) > 1000
    has_short_lines = sum(1 for line in text.splitlines() if len(line.strip()) < 80) / max(len(text.splitlines()), 1) > 0.7
    
    # Log content characteristics for debugging
    logger.info(f"Document analysis - Size: {text_size/1024:.1f}KB, Type: {ext}, "
                f"Markdown: {is_markdown}, Code: {is_code}, "
                f"Long paragraphs: {has_long_paragraphs}, Short lines: {has_short_lines}")
    
    # Dynamically determine chunk size based on content size and type - optimized for comprehensive coverage
    if text_size < 10000:  # Small documents (<10KB)
        # For small documents, use smaller chunks with high overlap to capture everything
        base_chunk_size = 250  
        overlap = 100  # 40% overlap to ensure nothing is missed
    elif text_size < 100000:  # Medium documents (<100KB)
        base_chunk_size = 400  # Smaller chunks for better coverage
        overlap = 200  # 50% overlap to ensure context continuity
    else:  # Large documents (>100KB)
        base_chunk_size = 600  # Moderate chunks for big documents
        overlap = 300  # 50% overlap to ensure context is maintained
    
    # Adjust chunk size based on content type
    if is_code:
        # Code needs smaller chunks for precision
        chunk_size = int(base_chunk_size * 0.8)
        if has_short_lines:
            # Code with many short lines needs even smaller chunks
            chunk_size = int(chunk_size * 0.8)
    elif is_markdown:
        # Markdown can use medium chunks
        chunk_size = base_chunk_size
    elif has_long_paragraphs:
        # Documents with long paragraphs need larger chunks
        chunk_size = int(base_chunk_size * 1.2)
    else:
        chunk_size = base_chunk_size
    
    # Select appropriate separators based on content type
    if is_markdown:
        # For markdown, use semantic chunking based on headers
        logger.info(f"Using markdown-optimized chunking strategy (chunk size: {chunk_size})")
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            length_function=len,
            separators=["\n## ", "\n# ", "\n### ", "\n\n", "\n", ". ", " ", ""]  # Prioritize markdown headers
        )
    elif is_code:
        # For code, use line-based chunking with function/class awareness
        logger.info(f"Using code-optimized chunking strategy (chunk size: {chunk_size})")
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            length_function=len,
            separators=["\nclass ", "\ndef ", "\nfunction ", "\n\n", "\n", " "]
        )
    else:
        # For other documents, use paragraph-based chunking
        logger.info(f"Using general document chunking strategy (chunk size: {chunk_size})")
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", ", ", " ", ""]
        )
    
    # Print a sample of the text for debugging
    logger.info(f"Text sample (first 100 chars): {text[:100]}...")
    
    # Extract document header (first 1000 characters) to preserve author information
    document_header = text[:1000]
    
    # Check if the header likely contains author information
    has_author_info = any(keyword in document_header.lower() for keyword in 
                         ["author", "authors", "correspondence", "affiliation", "university", "institute"])
    
    # Create chunks
    chunks = text_splitter.split_text(text)
    logger.info(f"Created {len(chunks)} chunks from document")
    
    # Create documents with metadata and add chunk titles for better context
    documents = []
    
    # If author information is detected, create a special "header" chunk
    if has_author_info:
        logger.info("Author information detected in document header - creating special header chunk")
        # Create a special chunk for the document header
        documents.append({
            "page_content": document_header,
            "metadata": {
                "source": filename,
                "chunk": "header",
                "title": "Document Header (Author Information)",
                "char_count": len(document_header),
                "is_header": True,
                "content_type": "author_info"
            }
        })
    
    # Process regular chunks
    for i, chunk in enumerate(chunks):
        # Extract a title from the chunk (first line or first few words)
        title = chunk.split('\n')[0][:30].strip()
        if not title:
            title = chunk[:30].strip()
        
        # Add chunk number to title
        chunk_title = f"Chunk {i+1}: {title}..."
        
        # Log the first chunk for debugging
        if i == 0:
            logger.info(f"First chunk sample: {chunk[:100]}...")
        
        # Determine if this chunk might contain author information
        chunk_lower = chunk.lower()
        is_author_chunk = any(keyword in chunk_lower for keyword in 
                             ["author", "authors", "correspondence", "affiliation", "university", "institute"])
        
        # Create document with enhanced metadata
        documents.append({
            "page_content": chunk,
            "metadata": {
                "source": filename,
                "chunk": i,
                "title": chunk_title,
                "char_count": len(chunk),
                "is_header": False,
                "content_type": "author_info" if is_author_chunk else "content"
            }
        })
    
    return documents

def process_document(file_path: str, original_filename: str):
    """Process a document: extract text, split into chunks, and create embeddings"""
    start_time = time.time()
    logger.info(f"Processing document: {original_filename}")
    
    # Extract text from document
    text = extract_text(file_path)
    
    if not text:
        logger.warning(f"No text extracted from {original_filename}")
        return
    
    # Create chunks
    chunks = create_document_chunks(text, original_filename)
    
    if not chunks:
        logger.warning(f"No chunks created for {original_filename}")
        return
    
    # Create vector store
    try:
        # Import Document class from langchain
        from langchain_core.documents import Document
        
        # Create proper Document objects
        documents = [
            Document(
                page_content=chunk["page_content"],
                metadata=chunk["metadata"]
            ) for chunk in chunks
        ]
        
        # Use the original filename (without extension) as the directory name
        filename_base = os.path.splitext(os.path.basename(original_filename))[0]
        save_path = os.path.join(EMBEDDINGS_DIR, filename_base)
        
        # Remove existing embeddings if they exist
        if os.path.exists(save_path):
            try:
                shutil.rmtree(save_path)
                logger.info(f"Removed existing embeddings at {save_path}")
            except Exception as e:
                logger.warning(f"Could not remove existing embeddings: {e}")
        
        # Create FAISS index
        logger.info(f"Creating vector store with {len(documents)} documents")
        vector_store = FAISS.from_documents(documents, embeddings)
        
        # Save the index
        logger.info(f"Saving embeddings to {save_path}")
        try:
            vector_store.save_local(save_path)
            logger.info(f"Successfully saved embeddings to {save_path}")
        except TypeError as e:
            if "allow_dangerous_deserialization" in str(e):
                logger.warning("Ignoring allow_dangerous_deserialization parameter error")
                # Try an alternative approach
                import pickle
                os.makedirs(save_path, exist_ok=True)
                with open(os.path.join(save_path, "index.pkl"), "wb") as f:
                    pickle.dump(vector_store, f)
                logger.info("Saved using pickle as fallback")
            else:
                raise
        
        processing_time = time.time() - start_time
        logger.info(f"Document processed successfully: {original_filename} in {processing_time:.2f} seconds")
    except Exception as e:
        logger.error(f"Error creating embeddings for {original_filename}: {e}")
        import traceback
        logger.error(traceback.format_exc())